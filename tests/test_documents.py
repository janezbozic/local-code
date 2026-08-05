#!/usr/bin/env python3
"""Integration fixtures for preserving PDF/DOCX/PPTX/XLSX imports."""

from __future__ import annotations

import json
import pathlib
import shutil
import subprocess
import tempfile
import unittest
import uuid

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter

ROOT = pathlib.Path(__file__).resolve().parents[1]


class DocumentLifecycleTests(unittest.TestCase):
    def setUp(self) -> None:
        self.work = pathlib.Path(tempfile.mkdtemp(prefix="local-workbench-doc-"))
        self.stem = f"fixture-{uuid.uuid4().hex}"
        self.created: list[pathlib.Path] = []

    def tearDown(self) -> None:
        shutil.rmtree(self.work, ignore_errors=True)
        for path in self.created:
            if path.is_dir():
                shutil.rmtree(path, ignore_errors=True)
            else:
                path.unlink(missing_ok=True)

    def make_fixtures(self) -> list[tuple[pathlib.Path, int]]:
        docx = self.work / f"{self.stem}-word.docx"
        document = Document()
        document.add_heading("Document fixture", 0)
        document.add_paragraph("Preserved DOCX text.")
        document.save(docx)

        pptx = self.work / f"{self.stem}-slides.pptx"
        deck = Presentation()
        deck.slides.add_slide(deck.slide_layouts[0]).shapes.title.text = "Slide fixture"
        deck.save(pptx)

        xlsx = self.work / f"{self.stem}-sheet.xlsx"
        book = Workbook()
        book.active["A1"] = "Sheet fixture"
        book.save(xlsx)

        pdf = self.work / f"{self.stem}-pages.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=200, height=200)
        with pdf.open("wb") as handle:
            writer.write(handle)
        return [(docx, 1), (pptx, 1), (xlsx, 1), (pdf, 1)]

    def test_representative_imports_and_manifests(self) -> None:
        for source, expected_count in self.make_fixtures():
            result = subprocess.run(
                [str(ROOT / ".venv/documents/bin/python"), str(ROOT / "tools/documents/workflow.py"), "import", str(source)],
                cwd=ROOT,
                text=True,
                capture_output=True,
            )
            self.assertEqual(result.returncode, 0, result.stderr)
            canonical = ROOT / "knowledge/markdown" / f"{source.stem}.md"
            originals = list((ROOT / "knowledge/originals").glob(f"*-{source.name}"))
            manifests = list((ROOT / "knowledge/manifests").glob(f"*-{source.stem}-*-import.json"))
            self.assertTrue(canonical.is_file())
            self.assertEqual(len(originals), 1)
            self.assertEqual(len(manifests), 1)
            record = json.loads(manifests[0].read_text())
            self.assertEqual(record["item_count"], expected_count)
            self.assertEqual(len(record["source_sha256"]), 64)
            self.created.extend([canonical, originals[0], manifests[0]])


if __name__ == "__main__":
    unittest.main()
